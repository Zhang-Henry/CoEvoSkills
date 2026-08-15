"""
Utility functions for detecting and processing dangling paper titles in PPTX files.
"""
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy


def detect_dangling_titles(prs):
    """
    Detect dangling paper titles in a presentation.
    These are TEXT_BOX shapes (not placeholders) that contain paper title text.
    Returns list of dicts with slide_index, shape, and title text.
    """
    dangling = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check if it's a text box (not a placeholder)
            is_placeholder = False
            try:
                if shape.placeholder_format is not None:
                    is_placeholder = True
            except ValueError:
                pass
            
            if not is_placeholder and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:  # Non-empty text box
                    dangling.append({
                        'slide_index': slide_idx,
                        'shape': shape,
                        'text': text,
                    })
    return dangling


def estimate_text_width_emu(text, font_size_pt=16, font_name='Arial'):
    """
    Estimate the width needed for text in EMU at given font size.
    Uses approximate character width for Arial font.
    Arial at 16pt: average character width ~9.5pt for mixed case text.
    We use a conservative estimate to ensure single-line display.
    """
    # For Arial 16pt, approximate average char width
    # Capital letters ~11pt, lowercase ~7.5pt, spaces ~4pt
    # Use a weighted average based on typical text
    total_width_pt = 0
    for ch in text:
        if ch == ' ':
            total_width_pt += font_size_pt * 0.35
        elif ch.isupper():
            total_width_pt += font_size_pt * 0.72
        elif ch in '.:,;-':
            total_width_pt += font_size_pt * 0.35
        elif ch in 'mwMW':
            total_width_pt += font_size_pt * 0.85
        elif ch in 'ilIjft':
            total_width_pt += font_size_pt * 0.35
        else:
            total_width_pt += font_size_pt * 0.58
    
    # Add margins/padding (left+right internal margins)
    total_width_pt += 20  # ~10pt padding each side
    
    # Convert pt to EMU (1 pt = 12700 EMU)
    width_emu = int(total_width_pt * 12700)
    return width_emu


def format_dangling_title(shape, slide_width, slide_height):
    """
    Format a dangling title shape:
    - Font: Arial, 16pt, color #989596, not bold
    - Adjust box width for single line
    - Position at bottom center of slide
    """
    text = shape.text_frame.text.strip()
    
    # Calculate needed width
    needed_width = estimate_text_width_emu(text, font_size_pt=16)
    
    # Ensure width doesn't exceed slide width (with margins)
    max_width = slide_width - Emu(200000)  # small margin
    if needed_width > max_width:
        needed_width = max_width
    
    # Set shape width
    shape.width = needed_width
    
    # Set height for single line of 16pt text + padding
    line_height = int(Pt(16) * 1.5)  # 1.5x line height
    shape.height = line_height + Emu(50000)  # Add some padding
    
    # Position at bottom center
    shape.left = (slide_width - shape.width) // 2
    # Place near bottom with some margin from bottom edge
    bottom_margin = Emu(200000)  # ~0.22 inches from bottom
    shape.top = slide_height - shape.height - bottom_margin
    
    # Format all runs in all paragraphs
    for para in shape.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x98, 0x95, 0x96)
            run.font.bold = False
    
    # Also set paragraph-level font defaults via XML to handle inherited properties
    for para in shape.text_frame.paragraphs:
        # If there are no runs but there is text, we need to handle that
        if not para.runs and para.text.strip():
            # Access the XML directly to set default run properties
            pPr = para._p.get_or_add_pPr()
            defRPr = pPr.find(qn('a:defRPr'))
            if defRPr is None:
                defRPr = pPr.makeelement(qn('a:defRPr'), {})
                pPr.append(defRPr)
            defRPr.set('sz', str(int(Pt(16))))
            defRPr.set('b', '0')
    
    # Set text frame properties - no word wrap, auto-size
    shape.text_frame.word_wrap = False
    
    # Set vertical anchor to middle
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def collect_unique_titles(dangling_items):
    """
    Collect unique paper titles from dangling items, preserving order of first appearance.
    """
    seen = set()
    unique_titles = []
    for item in dangling_items:
        title = item['text'].strip()
        if title not in seen:
            seen.add(title)
            unique_titles.append(title)
    return unique_titles


def create_reference_slide(prs, titles):
    """
    Create a new slide at the end with 'Reference' as title
    and all paper titles as auto-numbered bullet points.
    """
    # Try to find a layout with title and content
    layout = None
    for sl in prs.slide_layouts:
        if 'Title and Content' in sl.name:
            layout = sl
            break
    if layout is None:
        # Use the second layout (typically Title and Content)
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    
    slide = prs.slides.add_slide(layout)
    
    # Set the title
    title_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:  # Title placeholder
            title_shape = shape
            break
    
    if title_shape:
        title_shape.text = 'Reference'
    
    # Find the content/body placeholder
    body_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder
            body_shape = shape
            break
    
    if body_shape:
        tf = body_shape.text_frame
        tf.clear()
        
        for i, title in enumerate(titles):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            
            para.text = title
            para.level = 0
            
            # Set auto-numbered bullet points
            pPr = para._p.get_or_add_pPr()
            buNone = pPr.find(qn('a:buNone'))
            if buNone is not None:
                pPr.remove(buNone)
            
            # Remove any existing bullet format
            for child in list(pPr):
                if child.tag in [qn('a:buChar'), qn('a:buAutoNum'), qn('a:buNone'), qn('a:buBlip')]:
                    pPr.remove(child)
            
            # Add auto-numbered bullet
            buAutoNum = pPr.makeelement(qn('a:buAutoNum'), {'type': 'arabicPeriod'})
            pPr.append(buAutoNum)
    
    return slide


def process_pptx(input_path, output_path):
    """
    Main entry point: process the PPTX file.
    1. Detect dangling paper titles
    2. Reformat them (font, size, color, bold)
    3. Adjust width and position (bottom center)
    4. Create reference slide with unique titles
    5. Save to output path
    """
    prs = Presentation(input_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Step 1: Detect dangling titles
    dangling = detect_dangling_titles(prs)
    print(f"Found {len(dangling)} dangling paper titles:")
    for d in dangling:
        print(f"  Slide {d['slide_index']+1}: {d['text']}")
    
    # Step 2 & 3: Format and reposition each dangling title
    for d in dangling:
        format_dangling_title(d['shape'], slide_width, slide_height)
        print(f"  Formatted and repositioned: {d['text'][:60]}")
    
    # Step 4: Collect unique titles and create reference slide
    unique_titles = collect_unique_titles(dangling)
    print(f"\nUnique titles for reference slide ({len(unique_titles)}):")
    for i, t in enumerate(unique_titles):
        print(f"  {i+1}. {t}")
    
    create_reference_slide(prs, unique_titles)
    print("\nReference slide created.")
    
    # Step 5: Save
    prs.save(output_path)
    print(f"\nSaved to {output_path}")
    
    return {
        'dangling_count': len(dangling),
        'unique_titles': unique_titles,
        'output_path': output_path
    }


def validate_output(output_path):
    """
    Validate the processed PPTX file.
    """
    prs = Presentation(output_path)
    errors = []
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Check that dangling titles exist and are properly formatted
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            is_placeholder = False
            try:
                if shape.placeholder_format is not None:
                    is_placeholder = True
            except ValueError:
                pass
            
            if not is_placeholder and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and slide_idx < len(prs.slides) - 1:  # Not the reference slide
                    # Check formatting
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name != 'Arial':
                                errors.append(f"Slide {slide_idx+1}: font is {run.font.name}, expected Arial")
                            if run.font.size != Pt(16):
                                errors.append(f"Slide {slide_idx+1}: size is {run.font.size}, expected {Pt(16)}")
                            if run.font.bold:
                                errors.append(f"Slide {slide_idx+1}: bold is True, expected False")
                            try:
                                if run.font.color.rgb != RGBColor(0x98, 0x95, 0x96):
                                    errors.append(f"Slide {slide_idx+1}: color is {run.font.color.rgb}, expected 989596")
                            except:
                                errors.append(f"Slide {slide_idx+1}: could not check color")
                    
                    # Check position - should be bottom center
                    center_x = slide_width // 2
                    shape_center_x = shape.left + shape.width // 2
                    if abs(shape_center_x - center_x) > Emu(50000):
                        errors.append(f"Slide {slide_idx+1}: not horizontally centered")
    
    # Check reference slide (last slide)
    last_slide = prs.slides[-1] if len(prs.slides) > 0 else None
    if last_slide:
        has_reference_title = False
        has_numbered_list = False
        for shape in last_slide.placeholders:
            if shape.placeholder_format.idx == 0:
                if 'Reference' in shape.text:
                    has_reference_title = True
            if shape.placeholder_format.idx == 1:
                tf = shape.text_frame
                if len(tf.paragraphs) > 0:
                    has_numbered_list = True
        
        if not has_reference_title:
            errors.append("Reference slide: title does not contain 'Reference'")
        if not has_numbered_list:
            errors.append("Reference slide: no numbered list found")
    else:
        errors.append("No slides found in output")
    
    if errors:
        print("Validation errors:")
        for e in errors:
            print(f"  - {e}")
    else:
        print("Validation passed!")
    
    return errors
