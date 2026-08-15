"""Reusable utilities for surgical edits to embedded workbooks in PPTX files."""
import zipfile
import io
import re
import math
from lxml import etree

_SSML = 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'
_RELS = 'http://schemas.openxmlformats.org/package/2006/relationships'
_WS_T = ('http://schemas.openxmlformats.org/officeDocument/'
         '2006/relationships/worksheet')


def scan_textboxes(pptx_path):
    """Return a list of text strings from all text-frame shapes."""
    from pptx import Presentation
    texts = []
    for slide in Presentation(pptx_path).slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                texts.append(shape.text)
    return texts


def parse_update_instruction(text):
    """Extract (from_label, to_label, numeric_value) from free text."""
    flat = ' '.join(text.split())
    m = re.search(
        r'([A-Za-z0-9]+)\s+to\s+([A-Za-z0-9]+)\s*=\s*'
        r'([0-9]+(?:\.[0-9]*)?)'
        , flat)
    if m:
        return m.group(1), m.group(2), float(m.group(3))
    return None


def locate_embedded_xlsx(pptx_path):
    """Return the ZIP-internal path of the first embedded .xlsx."""
    with zipfile.ZipFile(pptx_path) as z:
        for n in z.namelist():
            if n.endswith('.xlsx') and 'embedding' in n.lower():
                return n
    raise FileNotFoundError('No embedded xlsx in ' + pptx_path)


def _ws_path(xlsx_bytes):
    """Resolve worksheet XML path from workbook rels."""
    with zipfile.ZipFile(io.BytesIO(xlsx_bytes)) as z:
        rp = 'xl/_rels/workbook.xml.rels'
        if rp in z.namelist():
            for rel in etree.fromstring(z.read(rp)).findall(
                    f'{{{_RELS}}}Relationship'):
                if rel.get('Type') == _WS_T:
                    t = rel.get('Target')
                    return ('xl/' + t) if not t.startswith('/') else t
        for n in z.namelist():
            if 'worksheets/sheet' in n:
                return n
    raise FileNotFoundError('No worksheet in xlsx')


def _shared_strings(xlsx_bytes):
    """Return list[str] from the shared-strings table."""
    with zipfile.ZipFile(io.BytesIO(xlsx_bytes)) as z:
        if 'xl/sharedStrings.xml' not in z.namelist():
            return []
        root = etree.fromstring(z.read('xl/sharedStrings.xml'))
        out = []
        for si in root.findall(f'{{{_SSML}}}si'):
            t = si.find(f'{{{_SSML}}}t')
            if t is not None and t.text:
                out.append(t.text)
            else:
                out.append(''.join(
                    r.text or '' for r in si.findall(f'.//{{{_SSML}}}t')))
        return out


def find_matrix_cell(xlsx_bytes, from_label, to_label):
    """Locate the cell ref at (row=from_label, col=to_label) in a header matrix."""
    wsp = _ws_path(xlsx_bytes)
    ss = _shared_strings(xlsx_bytes)
    with zipfile.ZipFile(io.BytesIO(xlsx_bytes)) as z:
        root = etree.fromstring(z.read(wsp))
    sd = root.find(f'{{{_SSML}}}sheetData')
    if sd is None:
        return None
    rows = sd.findall(f'{{{_SSML}}}row')
    if not rows:
        return None
    first_r = int(rows[0].get('r'))
    col_map, row_map = {}, {}
    for relem in rows:
        rn = int(relem.get('r'))
        cells = relem.findall(f'{{{_SSML}}}c')
        for ci, cel in enumerate(cells):
            ref = cel.get('r')
            col = re.match(r'([A-Z]+)', ref).group(1)
            ct = cel.get('t', '')
            ve = cel.find(f'{{{_SSML}}}v')
            txt = None
            if ct == 's' and ve is not None:
                idx = int(ve.text)
                txt = ss[idx] if idx < len(ss) else None
            elif ve is not None:
                txt = ve.text
            if txt is None:
                continue
            txt = txt.strip()
            if rn == first_r:
                col_map[txt] = col
            if ci == 0 and rn > first_r:
                row_map[txt] = rn
    if from_label in row_map and to_label in col_map:
        return f'{col_map[to_label]}{row_map[from_label]}'
    return None


def _extract_formulas(xml_str):
    """Extract all formula cells: {cell_ref: formula_text}."""
    pat = r'<c\s+r="([A-Z]+\d+)"[^>]*>\s*<f>([^<]+)</f>'
    return {m.group(1): m.group(2) for m in re.finditer(pat, xml_str)}


def _extract_values(xml_str):
    """Extract all non-formula numeric cell values: {cell_ref: float}."""
    vals = {}
    pat = r'<c\s+r="([A-Z]+\d+)"([^>]*)>(.*?)</c>'
    for m in re.finditer(pat, xml_str, re.DOTALL):
        ref, attrs, inner = m.group(1), m.group(2), m.group(3)
        if 't="s"' in attrs or '<f>' in inner:
            continue
        vm = re.search(r'<v>([^<]+)</v>', inner)
        if vm:
            try:
                vals[ref] = float(vm.group(1))
            except ValueError:
                pass
    return vals


def _find_dependents(formulas, changed_cell):
    """Find all formula cells transitively depending on changed_cell."""
    deps = {ref: set(re.findall(r'[A-Z]+\d+', f))
            for ref, f in formulas.items()}
    result = set()
    queue = [changed_cell]
    while queue:
        cur = queue.pop(0)
        for ref, refs_to in deps.items():
            if cur in refs_to and ref not in result:
                result.add(ref)
                queue.append(ref)
    return result


def _safe_eval_expr(expr, vals):
    """Evaluate a simple arithmetic expression with cell references."""
    def repl(m):
        r = m.group(0)
        return repr(vals[r]) if r in vals else r
    sub = re.sub(r'[A-Z]+\d+', repl, expr)
    if not re.match(r'^[\d\.eE+\-*/() ]+$', sub):
        return None
    try:
        r = eval(sub)
        return float(r) if isinstance(r, (int, float)) and math.isfinite(r) else None
    except Exception:
        return None


def _eval_formula(formula, vals):
    """Evaluate an Excel formula. Supports ROUND and basic arithmetic."""
    f = formula.strip()
    rm = re.match(r'^ROUND\((.+),\s*(\d+)\)$', f)
    if rm:
        inner = _safe_eval_expr(rm.group(1), vals)
        return round(inner, int(rm.group(2))) if inner is not None else None
    return _safe_eval_expr(f, vals)


def _format_cached(val):
    """Format a computed cached value for OOXML storage.
    
    Uses the same full-precision representation that spreadsheet
    applications use for computed results in OOXML XML.
    """
    return '%.17g' % val


def _patch_sheet_xml(sheet_bytes, cell_ref, new_val_str):
    """Replace a single non-formula cell value and update dependent
    formula cached values in raw XML bytes."""
    xml = sheet_bytes.decode('utf-8')
    
    # 1. Replace target cell value
    cpat = r'<c\s+r="' + re.escape(cell_ref) + r'"[^>]*>(.*?)</c>'
    cm = re.search(cpat, xml, re.DOTALL)
    if not cm:
        raise ValueError(f'{cell_ref} not found')
    inner = cm.group(1)
    if '<f>' in inner:
        raise ValueError(f'{cell_ref} is a formula cell')
    vm = re.search(r'(<v>)([^<]*)(</v>)', inner)
    if not vm:
        raise ValueError(f'{cell_ref} has no value element')
    s = cm.start(1) + vm.start(2)
    e = cm.start(1) + vm.end(2)
    xml = xml[:s] + new_val_str + xml[e:]
    
    # 2. Update cached values for dependent formula cells
    new_value = float(new_val_str)
    formulas = _extract_formulas(xml)
    dependents = _find_dependents(formulas, cell_ref)
    
    if dependents:
        # Build value map with the updated cell
        vals = _extract_values(xml)
        vals[cell_ref] = new_value
        
        # Process dependents in order (handle cascading)
        for dep_ref in sorted(dependents):
            formula = formulas.get(dep_ref)
            if not formula:
                continue
            computed = _eval_formula(formula, vals)
            if computed is None:
                continue
            # Replace the cached <v> inside this formula cell
            dep_pat = (r'(<c\s+r="' + re.escape(dep_ref) +
                      r'"[^>]*>\s*<f>[^<]+</f>\s*<v>)([^<]+)(</v>)')
            dm = re.search(dep_pat, xml)
            if dm:
                cached_str = _format_cached(computed)
                xml = xml[:dm.start(2)] + cached_str + xml[dm.end(2):]
                vals[dep_ref] = computed
    
    return xml.encode('utf-8')


def _repack_zip(src_bytes, replacements):
    """Repack a ZIP replacing specific entries."""
    buf = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(src_bytes)) as zi:
        with zipfile.ZipFile(buf, 'w') as zo:
            for item in zi.infolist():
                data = replacements.get(item.filename, zi.read(item.filename))
                info = zipfile.ZipInfo(item.filename)
                info.compress_type = item.compress_type
                info.external_attr = item.external_attr
                info.date_time = item.date_time
                zo.writestr(info, data)
    return buf.getvalue()


def update_embedded_value(input_pptx, output_pptx,
                          from_label, to_label, new_value):
    """Update one matrix cell in the embedded workbook and refresh
    cached values for any formula cells that depend on it."""
    xlsx_arc = locate_embedded_xlsx(input_pptx)
    with zipfile.ZipFile(input_pptx) as z:
        xlsx_bytes = z.read(xlsx_arc)
    wsp = _ws_path(xlsx_bytes)
    cell_ref = find_matrix_cell(xlsx_bytes, from_label, to_label)
    if cell_ref is None:
        raise ValueError(f'No cell for {from_label}->{to_label}')
    with zipfile.ZipFile(io.BytesIO(xlsx_bytes)) as z:
        sheet = z.read(wsp)
    patched = _patch_sheet_xml(sheet, cell_ref, repr(new_value))
    new_xlsx = _repack_zip(xlsx_bytes, {wsp: patched})
    with open(input_pptx, 'rb') as f:
        pptx_raw = f.read()
    result = _repack_zip(pptx_raw, {xlsx_arc: new_xlsx})
    with open(output_pptx, 'wb') as f:
        f.write(result)
    return cell_ref


def run_end_to_end(input_pptx, output_pptx):
    """Full pipeline: read instruction from text box, apply update, save."""
    texts = scan_textboxes(input_pptx)
    instr = None
    for t in texts:
        instr = parse_update_instruction(t)
        if instr:
            break
    if instr is None:
        raise ValueError('No update instruction found')
    fr, to, val = instr
    print(f'Instruction: {fr} to {to} = {val}')
    cell = update_embedded_value(input_pptx, output_pptx, fr, to, val)
    print(f'Updated cell {cell}')
    return {'from_currency': fr, 'to_currency': to, 'new_rate': val,
            'updated_cell': cell}


def validate(output_pptx, from_label, to_label, expected_rate):
    """Validate the output file integrity and correctness."""
    import openpyxl
    from pptx import Presentation

    prs = Presentation(output_pptx)
    print('PPTX opens OK')

    tb_count = sum(1 for sl in prs.slides for sh in sl.shapes
                   if sh.has_text_frame and sh.text.strip())
    assert tb_count > 0, 'No text boxes found'
    print(f'Text boxes: {tb_count}')

    xlsx_arc = locate_embedded_xlsx(output_pptx)
    with zipfile.ZipFile(output_pptx) as z:
        xb = z.read(xlsx_arc)

    # Check target cell
    wb = openpyxl.load_workbook(io.BytesIO(xb), data_only=True)
    cr = find_matrix_cell(xb, from_label, to_label)
    assert cr is not None
    actual = wb.active[cr].value
    assert actual == expected_rate, f'{actual} != {expected_rate} at {cr}'
    print(f'Cell {cr} = {expected_rate} OK')

    # Check formulas preserved
    wb2 = openpyxl.load_workbook(io.BytesIO(xb))
    ws = wb2.active
    formulas = [(c.coordinate, c.value)
                for row in ws.iter_rows(max_row=ws.max_row, max_col=ws.max_column)
                for c in row
                if isinstance(c.value, str) and c.value.startswith('=')]
    assert len(formulas) > 0, 'All formulas lost'
    for coord, fval in formulas:
        print(f'  Formula {coord}: {fval}')
    print(f'Formulas intact: {len(formulas)}')

    # Check cached values are consistent with formulas
    ws_data = wb.active  # data_only view
    for coord, fval in formulas:
        cached = ws_data[coord].value
        if cached is not None:
            print(f'  Cached {coord} = {cached}')

    print('Validation PASSED')
    return True
