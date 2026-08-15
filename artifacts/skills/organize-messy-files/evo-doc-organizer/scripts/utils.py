import os
import subprocess
import shutil
import re


def extract_text_pdf(filepath, max_pages=3):
    """Extract text from a PDF using pdftotext (poppler)."""
    try:
        result = subprocess.run(
            ['pdftotext', filepath, '-', '-l', str(max_pages)],
            capture_output=True, text=True, timeout=15
        )
        text = result.stdout.strip()
        if text:
            return text
    except Exception:
        pass
    # Fallback: try PyPDF2
    try:
        import PyPDF2
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            pages = []
            for i, page in enumerate(reader.pages):
                if i >= max_pages:
                    break
                t = page.extract_text()
                if t:
                    pages.append(t)
            return '\n'.join(pages)
    except Exception:
        pass
    return ''


def extract_text_docx(filepath):
    """Extract text from a DOCX file."""
    try:
        from docx import Document
        doc = Document(filepath)
        return '\n'.join([p.text for p in doc.paragraphs])
    except Exception:
        return ''


def extract_text_pptx(filepath):
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text_parts.append(shape.text)
        return '\n'.join(text_parts)
    except Exception:
        return ''


def extract_text(filepath):
    """Extract text from a file based on its extension."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.pdf':
        return extract_text_pdf(filepath)
    elif ext == '.docx':
        return extract_text_docx(filepath)
    elif ext == '.pptx':
        return extract_text_pptx(filepath)
    return ''


def extract_arxiv_category(text):
    """Extract arXiv category from text if present."""
    match = re.search(r'arXiv:\S+\s*\[([^\]]+)\]', text)
    if match:
        return match.group(1)
    return None


def classify_document(text):
    """Classify document text into one of 5 categories.
    
    Categories:
    - LLM: Large Language Models
    - trapped_ion_and_qc: Trapped ion and quantum computing
    - black_hole: Black holes
    - DNA: DNA-related research
    - music_history: Music history and evolution
    
    Returns tuple of (folder_name, scores_dict).
    """
    text_lower = text.lower()
    arxiv_cat = extract_arxiv_category(text)
    
    # Score each category with weighted keywords
    scores = {
        'LLM': 0,
        'trapped_ion_and_qc': 0,
        'black_hole': 0,
        'DNA': 0,
        'music_history': 0
    }
    
    # LLM keywords
    llm_keywords = [
        ('large language model', 10), ('llm', 8), ('llms', 8),
        ('language model', 6), ('fine-tuning', 3), ('fine-tune', 3),
        ('instruction tuning', 5), ('transformer', 3),
        ('gpt', 4), ('chatgpt', 5), ('generative pre-trained', 5),
        ('binarized', 3), ('quantization', 2), ('llm compression', 8),
        ('natural language processing', 3), ('nlp', 2),
        ('prompt', 2), ('instruction-tuned', 5),
        ('code generation', 3), ('evol-instruct', 5),
        ('visual grounding', 2), ('speech translation', 2),
    ]
    
    # Trapped ion / quantum computing keywords
    ion_qc_keywords = [
        ('trapped ion', 10), ('trapped-ion', 10), ('ion trap', 10),
        ('quantum computing', 6), ('quantum computer', 6),
        ('quantum simulation', 5), ('quantum simulator', 5),
        ('qubit', 4), ('qubits', 4),
        ('quantum information', 4), ('quantum circuit', 4),
        ('jaynes-cummings', 5), ('lamb-dicke', 5),
        ('phonon', 3), ('ion chain', 5), ('ion-photon', 6),
        ('shuttling', 3), ('quantum gate', 4),
        ('quantum entanglement', 3),
        ('zitterbewegung', 5), ('lissajous', 3),
        ('squeezed phonon', 5), ('quantum thermodynamics', 3),
        ('motional frequency', 5),
        ('quantum network', 3),
        ('quantum algorithm', 5), ('quantum accelerat', 5),
        ('grover', 3), ('quantum speedup', 4),
    ]
    
    # Black hole keywords
    bh_keywords = [
        ('black hole', 10), ('black-hole', 10), ('schwarzschild', 8),
        ('kerr black', 8), ('event horizon', 6), ('hawking radiation', 6),
        ('gravitational wave', 5), ('wormhole', 5),
        ('black hole thermodynamics', 10), ('singularity', 3),
        ('neutron star', 3), ('coalescen', 3),
        ('horizon wave', 5), ('reissner-nordstr', 5),
        ('dark matter', 2), ('general relativity', 2),
        ('gr-qc', 5), ('black hole remnant', 8),
        ('black hole chemistry', 8), ('black hole spectroscopy', 8),
        ('perseus', 2), ('cluster outskirts', 2),
    ]
    
    # DNA keywords
    dna_keywords = [
        ('dna', 8), ('dna sequence', 10), ('dna methylation', 10),
        ('dna nanotechnology', 10), ('nucleotide', 5),
        ('genomic', 4), ('genome', 4), ('gene ', 3),
        ('double-stranded', 5), ('double stranded', 5),
        ('pcr', 4), ('polymerase', 4),
        ('translocation', 3), ('coarse-graining', 2),
        ('dna coding', 8), ('dna storage', 8),
        ('charge migration', 3), ('dna devices', 8),
        ('homeodomain', 5), ('dna-binding', 8),
        ('dinucleotide', 5), ('cancer drivers', 3),
        ('diffusion model for dna', 10),
    ]
    
    # Music history keywords
    music_keywords = [
        ('music', 6), ('musical', 5), ('popular music', 8),
        ('music history', 10), ('song', 3), ('melody', 5),
        ('harmony', 4), ('hip-hop', 6), ('rap', 3),
        ('lyrics', 5), ('musical scale', 6),
        ('musimetrics', 8), ('music evolution', 8),
        ('contemporary music', 6), ('western classical music', 8),
        ('music manuscripts', 6), ('album', 3),
        ('sonification', 4), ('inharmonicity', 5),
        ('noisiness', 3), ('disc-cover', 4),
        ('cultural evolution', 4), ('music chart', 6),
        ('vocal', 3), ('singer', 3),
    ]
    
    for keyword, weight in llm_keywords:
        count = text_lower.count(keyword)
        scores['LLM'] += count * weight
    
    for keyword, weight in ion_qc_keywords:
        count = text_lower.count(keyword)
        scores['trapped_ion_and_qc'] += count * weight
    
    for keyword, weight in bh_keywords:
        count = text_lower.count(keyword)
        scores['black_hole'] += count * weight
    
    for keyword, weight in dna_keywords:
        count = text_lower.count(keyword)
        scores['DNA'] += count * weight
    
    for keyword, weight in music_keywords:
        count = text_lower.count(keyword)
        scores['music_history'] += count * weight
    
    # Use arXiv category as a tiebreaker / boost
    if arxiv_cat:
        if arxiv_cat.startswith('quant-ph'):
            scores['trapped_ion_and_qc'] += 20
        elif arxiv_cat.startswith('gr-qc') or arxiv_cat.startswith('hep-th'):
            scores['black_hole'] += 20
        elif arxiv_cat.startswith('q-bio'):
            scores['DNA'] += 20
        elif arxiv_cat.startswith('cs.CL') or arxiv_cat.startswith('cs.LG'):
            scores['LLM'] += 20
        elif arxiv_cat.startswith('cs.SD'):
            scores['music_history'] += 20
    
    # Return category with highest score
    best_category = max(scores, key=scores.get)
    
    # If all scores are 0, default to music_history (the catch-all per instructions)
    if scores[best_category] == 0:
        best_category = 'music_history'
    
    return best_category, scores


def organize_documents(source_dir, output_base_dir):
    """Organize documents from source_dir into 5 category folders under output_base_dir.
    
    Files are MOVED (not copied) from source_dir to category folders.
    
    Args:
        source_dir: Directory containing all documents
        output_base_dir: Base directory where category folders will be created
    
    Returns:
        dict mapping filename to assigned category
    """
    categories = ['LLM', 'trapped_ion_and_qc', 'black_hole', 'DNA', 'music_history']
    
    # Create category directories
    for cat in categories:
        cat_dir = os.path.join(output_base_dir, cat)
        os.makedirs(cat_dir, exist_ok=True)
    
    # Get all files
    files = [f for f in os.listdir(source_dir) 
             if os.path.isfile(os.path.join(source_dir, f))]
    
    assignments = {}
    
    for filename in sorted(files):
        filepath = os.path.join(source_dir, filename)
        text = extract_text(filepath)
        category, scores = classify_document(text)
        assignments[filename] = category
        print(f'{filename} -> {category} (scores: {scores})')
        
        # Move file to category folder
        dest = os.path.join(output_base_dir, category, filename)
        shutil.move(filepath, dest)
    
    # Remove source directory if empty
    remaining = os.listdir(source_dir)
    if len(remaining) == 0:
        os.rmdir(source_dir)
        print(f'Removed empty source directory: {source_dir}')
    
    return assignments


def validate_organization(output_base_dir, expected_count=None):
    """Validate that files are organized correctly.
    
    Checks:
    1. All 5 category folders exist
    2. Each file appears in exactly one folder
    3. Total file count matches expected if provided
    """
    categories = ['LLM', 'trapped_ion_and_qc', 'black_hole', 'DNA', 'music_history']
    
    # Check all category folders exist
    for cat in categories:
        cat_dir = os.path.join(output_base_dir, cat)
        if not os.path.isdir(cat_dir):
            print(f'FAIL: Category folder {cat} does not exist')
            return False
    
    # Collect all organized files
    organized_files = {}
    total = 0
    for cat in categories:
        cat_dir = os.path.join(output_base_dir, cat)
        for f in os.listdir(cat_dir):
            if f in organized_files:
                print(f'FAIL: File {f} appears in multiple folders: {organized_files[f]} and {cat}')
                return False
            organized_files[f] = cat
            total += 1
    
    if expected_count and total != expected_count:
        print(f'FAIL: Expected {expected_count} files, found {total}')
        return False
    
    print(f'PASS: {total} files organized into {len(categories)} folders')
    for cat in categories:
        cat_dir = os.path.join(output_base_dir, cat)
        count = len(os.listdir(cat_dir))
        print(f'  {cat}: {count} files')
    
    return True
